from django.shortcuts import render
from django.http import HttpResponse, JsonResponse, FileResponse, Http404
from django.views.decorators.csrf import csrf_exempt
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough
from datetime import datetime
import xml.etree.ElementTree as ET
from pptx.dml.color import RGBColor
from pptx import Presentation
from docx import Document
import os
import subprocess

# Create your views here.

PRESENTATION_SYSTEM_PROMPT = """You are an agent whose task is to generate text in order to generate a PowerPoint presentation. 

In order to do this, you need generate text in an XML style format.

You will be using preset slide layouts. Each of these slide layouts has a list of placeholders that you will populate with text.

Here is an example of the format, showing each of the 6 preset slide layouts:

<slides>
    <slide layout="title">
        <title>Slide Title</title>
        <subtitle>Slide Subtitle</subtitle>
    </slide>
    <slide layout="content">
        <title>Slide Title</title>
        <content>
			<b>Point 1</b>
			<b>Point 2</b>
		</content>
    </slide>
    <slide layout="header">
	    <title>Slide Title</title>
	    <text>Text</text>
    </slide>
    <slide layout="two">
	    <title>Slide Title</title>
	    <content>
			<b>Point 1</b>
			<b>Point 2</b>
		</content>
	    <content>
			<b>Point 1</b>
			<b>Point 2</b>
	    </content>
    </slide>
    <slide layout="comp">
	    <title>Slide Title</title>
	    <text>Left Side Title</text>
	    <content>
			<b>Point 1</b>
			<b>Point 2</b>
	    </content>
	    <text>Right Side Title</text>
	    <content>
			<b>Point 1</b>
			<b>Point 2</b>
	    </content>
    </slide>
    <slide layout="caption">
	    <title>Slide Title</title>
	    <content>
			<b>Point 1</b>
			<b>Point 2</b>
	    </content>
	    <text>Caption</text>
    </slide>
</slides>

Please note that you do not need to use every single layout preset, only use the preset you determine to be the most effective at displaying the information for a particular slide.

Here is a description of each of the slide layouts for you to determine their usefulness:
"title": A slide with a centered title and optional subtitle
"content": A slide with a title on top and a content body below
"header": Shorthand for section header. A slide used to signal a new section in the presentation, essentially the same as a title slide but styled differently
"two": Short for two-content, a slide with a title on top, and two side by side content bodies
"comp": Short for comparison, the same thing as a two content slide, but with a title for each side
"caption": A slide with a title and caption on the left side, and a main content body on the right

Please also note that a content placeholder displays text as bullet points for every new line, whereas title/subtitle/text do not, so you will need to wrap each bullet point with a `<b>` tag.

Here is an example of what you need to produce:
<slides>
    <slide layout="title">
        <title>Understanding Subject</title>
        <subtitle>Looking at Subject</subtitle>
    </slide>
    <slide layout="content">
        <title>What is Topic?</title>
        <content>
			<b>Point 1</b>
			<b>Point 2</b>
		</content>
    </slide>
    <slide layout="content">
        <title>Why is topic Important?</title>
        <content>
			<b>Helps us learn new things</b>
		</content>
    </slide>
    <slide layout="comp">
	    <title>Topic</title>
	    <text>What is Topic?</text>
	    <content>
			<b>Point 1</b>
			<b>Point 2</b>
	    </content>
	    <text>Why is it ___?</text>
	    <content>
			<b>Point 1</b>
			<b>Point 2</b>
	    </content>
    </slide>
    <slide layout="content">
        <title>What is Topic?</title>
        <content>
			<b>Point 1</b>
            <b>Example: </b>
		</content>
    </slide>
    <slide layout="title">
        <title>Recap: </title>
        <subtitle>Understanding</subtitle>
    </slide>
</slides>

Please leave out any string used in markdown e.g. ```xml```
"""

PRESENTATION_GENERATE_TEMPLATE = """{system}
Based on the following text parsed from {ctx}. Help adapt this into new materials for a presentation can be used to assist learning for other age groups.

{document}

Please use the original materials and convert them into slides that can be can be used to teach a grade {targetGrade} class. {prompt}."""

PRESENTATION_REGENERATE_TEMPLATE = """{system}

This is text parsed from {ctx}. 
{document}

This is a previous presentation generated by an agent.
{originalString}

Here is what the user wants changed:
{prompt}

Please take the previous iteration, and modify it to satisfy the user's request. Refer to the original document when needed. 
"""

QUIZ_SYSTEM_PROMPT = """You are an agent whose task is to generate text in order to generate quizzes. 

In order to do this, you will need to output in a consistent XML style format. 
Please come up with a title, as well as multiple choice or short answer (depending on what the user asks for) questions for the quiz.  

Here is an example of the format:
<quiz>
	<title>Quiz Title</title>
	
	<question type="choice" ask="1. Question 1">
		<choice>A) Choice 1</choice>
		<choice>B) Choice 2</choice>
		<choice answer="true">C) Choice 3</choice>
		<choice>Choice 4</choice>
	</question>
	
	<question type="short">2. Short answer question 1</question>
</quiz>

Here is an explanation of the format:
The title tag defines the title of the quiz, please always start with a title tag. 
The question tag defines a question. A question has two types: "choice" and "short", referring to multiple choice and short answer questions respectively. For multiple choice questions, please provide the question in the "ask" attribute as shown in the example, and then provide a list of the choices nested in the question element as described in the example. For short answer questions, there is no "ask" attribute, instead, simply enter the question as the text of the question element as shown in the example.  
The choice tag defines a multiple choice question choice. A choice tag is given the attribute answer="true" if it is the correct answer to the question. 

Please note that this format was just an example to show how multiple choice and short answer questions are defined in the format, and that you may have as many questions as you deem necessary, but make at least 5 questions, in the case of more material up to 10 questions may be necessary. 
For the types of questions you use in your quiz, the user may request only multiple choice questions, only short answer questions, or a mixture of both multiple choice and short answer questions. 
In the case of the last option, a mixture of both multiple choice and short answer questions, please divide the quiz into two sections: multiple choice and short answer questions. 
In order to do this, start by asking all the multiple choice questions, and then asking the short answer questions, as shown in the following example:

<quiz>
	<title>Quiz Title</title>
	
	<question type="choice" ask="1. Question 1">
		<choice>A) Choice 1</choice>
		<choice>B) Choice 2</choice>
		<choice answer="true">C) Choice 3</choice>
		<choice>Choice 4</choice>
	</question>
 
    <question type="choice" ask="2. Question 2">
		<choice answer="true">A) Choice 1</choice>
		<choice>B) Choice 2</choice>
		<choice>C) Choice 3</choice>
		<choice>Choice 4</choice>
	</question>
 
    <question type="choice" ask="3. Question 3">
		<choice answer="true">A) Choice 1</choice>
		<choice>B) Choice 2</choice>
		<choice>C) Choice 3</choice>
		<choice>Choice 4</choice>
	</question>
	
	<question type="short">4. Short answer question 1</question>
    
    <question type="short">5. Short answer question 2</question>
</quiz>
"""

QUIZ_GENERATE_TEMPLATE = """{system}

Please generate a quiz based on the follwing text parsed from {ctx}. 
{document}

{targetGrade}

{questionType}

{prompt}
"""

QUIZ_REGENERATE_TEMPLATE = """{system}

This is a previous quiz generated by an agent. 
{originalString}

The following is text parsed from {ctx}, what the quiz is based on. 
{document}

Please take the previous iteration of the quiz, and modify it to satisfy the user's request. Referring to the original document if needed.
 
Here is what the user wants changed:
{prompt}

{questionType}
"""

current_file_path = os.path.dirname(os.path.realpath(__file__))
GENERATEDCONTENT_DIRECTORY = os.path.join(current_file_path, 'generatedcontent')

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
LIGHTBLUE = RGBColor(135, 206, 235)
CREAM = RGBColor(255, 253, 208)
GREY = RGBColor(128, 128, 128)

def generate_filename(extension, type, username):
    timestamp = datetime.now().strftime('%Y%m%d')
    file_name = f"{username}_{type}_{timestamp}{extension}"
    
    # handle duplicate filename
    counter = 1
    unique_file_name = file_name
    while os.path.exists(os.path.join(GENERATEDCONTENT_DIRECTORY, unique_file_name)):
        unique_file_name = f"{username}_{type}_{timestamp}_{counter}{extension}"
        counter += 1
    return unique_file_name
        
def set_background_color(prs, color):
    for master in prs.slide_masters:
        background = master.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color  # Set color to ---- (needs a form value to change here)

def set_font(prs, type):
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = type  # Set font to ---- (needs a form value to change here)

def set_font_color(prs, color):
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = color  # Set font color ---- (needs a form value to change)

def generate_slides_from_XML(xml_string, bgcolor, fonttype, fontcolor, username):
    root = ET.fromstring(xml_string)
    prs = Presentation()
    set_background_color(prs, bgcolor)

    for slide in root.findall('slide'):
        if slide.get('layout') == 'title':
            title_slide_layout = prs.slide_layouts[0]
            title_slide = prs.slides.add_slide(title_slide_layout)
            shapes = title_slide.shapes
            for i, element in enumerate(slide):
                if element.tag == "content":
                    content_string = "" 
                    for bullet in element:
                        content_string += (bullet.text + '\n')
                    shapes.placeholders[i].text_frame.text = content_string
                else:
                    shapes.placeholders[i].text_frame.text = element.text
        elif slide.get('layout') == 'content':
            content_slide_layout = prs.slide_layouts[1]
            content_slide = prs.slides.add_slide(content_slide_layout)
            shapes = content_slide.shapes
            for i, element in enumerate(slide):
                if element.tag == "content":
                    content_string = "" 
                    for bullet in element:
                        content_string += (bullet.text + '\n')
                    shapes.placeholders[i].text_frame.text = content_string
                else:
                    shapes.placeholders[i].text_frame.text = element.text
        elif slide.get('layout') == 'header':
            header_slide_layout = prs.slide_layouts[2]
            header_slide = prs.slides.add_slide(header_slide_layout)
            shapes = header_slide.shapes
            for i, element in enumerate(slide):
                if element.tag == "content":
                    content_string = "" 
                    for bullet in element:
                        content_string += (bullet.text + '\n')
                    shapes.placeholders[i].text_frame.text = content_string
                else:
                    shapes.placeholders[i].text_frame.text = element.text
        elif slide.get('layout') == 'two':
            twocontent_slide_layout = prs.slide_layouts[3]
            twocontent_slide = prs.slides.add_slide(twocontent_slide_layout)
            shapes = twocontent_slide.shapes
            for i, element in enumerate(slide):
                if element.tag == "content":
                    content_string = "" 
                    for bullet in element:
                        content_string += (bullet.text + '\n')
                    shapes.placeholders[i].text_frame.text = content_string
                else:
                    shapes.placeholders[i].text_frame.text = element.text
        elif slide.get('layout') == 'comp':
            comparison_slide_layout = prs.slide_layouts[4]
            comparison_slide = prs.slides.add_slide(comparison_slide_layout)
            shapes = comparison_slide.shapes
            for i, element in enumerate(slide):
                if element.tag == "content":
                    content_string = "" 
                    for bullet in element:
                        content_string += (bullet.text + '\n')
                    shapes.placeholders[i].text_frame.text = content_string
                else:
                    shapes.placeholders[i].text_frame.text = element.text
        elif slide.get('layout') == 'caption':
            caption_slide_layout = prs.slide_layouts[7]
            caption_slide = prs.slides.add_slide(caption_slide_layout)
            shapes = caption_slide.shapes
            for i, element in enumerate(slide):
                if element.tag == "content":
                    content_string = "" 
                    for bullet in element:
                        content_string += (bullet.text + '\n')
                    shapes.placeholders[i].text_frame.text = content_string
                else:
                    shapes.placeholders[i].text_frame.text = element.text
    set_font(prs, fonttype)
    set_font_color(prs, fontcolor)
    file_name_pptx = generate_filename(".pptx", "presentation", username)
    file_path_pptx = os.path.join(GENERATEDCONTENT_DIRECTORY, file_name_pptx)
    prs.save(file_path_pptx)
    
    file_name_pdf = generate_filename(".pdf", "presentation", username)
    convert_to_pdf(file_path_pptx, GENERATEDCONTENT_DIRECTORY)
    
    return file_name_pptx

def generate_quiz_from_XML(xml_string, username):
    document = Document()
    
    root = ET.fromstring(xml_string)
    title = root.find('title').text
    document.add_heading(title, 0)
    
    multiple_choice_answers = ""
    
    for question in root.findall('question'):
        if question.get('type') == 'short':
            document.add_paragraph().add_run(question.text).bold = True
            document.add_paragraph()
        if question.get('type') == 'choice':
            question_text = question.get('ask')
            document.add_paragraph().add_run(question_text).bold = True
            for choice in question.findall('choice'):
                document.add_paragraph(choice.text, style='List')
                if choice.get('answer') == 'true':
                    answer = choice.text[0] + " "
                    multiple_choice_answers += answer
    
    document.add_page_break()
    document.add_paragraph("Multiple Choice Answers: " + multiple_choice_answers)       
    # todo: save file
    file_name_docx = generate_filename(".docx", "quiz", username)
    file_path_docx = os.path.join(GENERATEDCONTENT_DIRECTORY, file_name_docx)
    document.save(file_path_docx)
    
    file_name_pdf = generate_filename(".pdf", "quiz", username)
    convert_to_pdf(file_path_docx, GENERATEDCONTENT_DIRECTORY)
    
    return file_name_docx

@csrf_exempt
def generate_presentation(request):
    if (request.method == 'POST'):
        load_dotenv()
        uploaded_file = request.FILES.get('file')
        user_prompt = request.POST.get('prompt')
        ctx = request.POST.get('ctx')
        targetGrade = request.POST.get('targetGrade')
        bgcolor = request.POST.get('backgroundColor')
        fonttype = request.POST.get('fontType')
        fontcolor = request.POST.get('fontColor')
        username = request.POST.get('username')

        # extract text from uploaded file
        text = ""
        
        if uploaded_file.name.endswith('.pdf'):
            pdf_reader = PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                text += page.extract_text()
        elif uploaded_file.name.endswith('.txt'):
            text += uploaded_file.read().decode("utf-8")
        elif uploaded_file.name.endswith('.docx'):
            doc = Document(uploaded_file)
            for paragraph in doc.paragraphs:
                text += paragraph.text + '\n'
        elif uploaded_file.name.endswith('.pptx'):
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + '\n'           
        else:
            return JsonResponse({'error': 'Unsupported file format'})
        
        # Split text into chunks
        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
		
        chunks = text_splitter.split_text(text)
        embeddings = OpenAIEmbeddings()
		
        vectorstore = FAISS.from_texts(chunks, embeddings)
		
        if ctx:
            model = ChatOpenAI(model_name="gpt-3.5-turbo-0125", temperature=0.8)
            filled_template = PRESENTATION_GENERATE_TEMPLATE.replace("{system}", PRESENTATION_SYSTEM_PROMPT)\
                                      .replace("{targetGrade}", targetGrade)\
                                      .replace("{prompt}", user_prompt)
            prompt = ChatPromptTemplate.from_template(template=filled_template)
            retriever = vectorstore.as_retriever(search_type="similarity_score_threshold", search_kwargs={"score_threshold": 0.5, "k": 4})
            
            chain = (
                {"document": retriever, "ctx": RunnablePassthrough()}
                | prompt
                | model
                | StrOutputParser()
            )
            response = chain.invoke(ctx)
            
            apply_bgcolor = WHITE
            if bgcolor == 'black':
                apply_bgcolor = BLACK
            elif bgcolor == 'white':
                apply_bgcolor = WHITE
            elif bgcolor == 'lightblue':
                apply_bgcolor = LIGHTBLUE
            elif bgcolor == 'cream':
                apply_bgcolor = CREAM
            elif bgcolor == 'grey':
                apply_bgcolor = GREY

            apply_fontcolor = BLACK
            if fontcolor == 'black':
                apply_fontcolor = BLACK
            elif fontcolor == 'white':
                apply_fontcolor = WHITE

            file_name = generate_slides_from_XML(response, apply_bgcolor, fonttype, apply_fontcolor, username)
            return JsonResponse({'filename' : file_name, 'response': response, 'file_text': text, 'context': ctx,'style': {'bg': bgcolor, 'fontcolor': fontcolor, 'fonttype': fonttype}})
        else:
            return JsonResponse({'response': 'No context specified'})
    return HttpResponse("Listening for requests...")

@csrf_exempt
def generate_quiz(request):
    if (request.method == 'POST'):
        load_dotenv()
        uploaded_file = request.FILES.get('file')
        user_prompt = request.POST.get('prompt')
        ctx = request.POST.get('ctx')
        targetGrade = request.POST.get('targetGrade')
        username = request.POST.get('username')
        question_type = request.POST.get('questionType')
        
        # extract text
        text = ""
        
        if uploaded_file.name.endswith('.pdf'):
            pdf_reader = PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                text += page.extract_text()
        elif uploaded_file.name.endswith('.txt'):
            text += uploaded_file.read().decode("utf-8")
        elif uploaded_file.name.endswith('.docx'):
            doc = Document(uploaded_file)
            for paragraph in doc.paragraphs:
                text += paragraph.text + '\n'
        elif uploaded_file.name.endswith('.pptx'):
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + '\n'           
        else:
            return JsonResponse({'error': 'Unsupported file format'})
        
         # Split text into chunks
        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
		
        chunks = text_splitter.split_text(text)
        embeddings = OpenAIEmbeddings()
		
        vectorstore = FAISS.from_texts(chunks, embeddings)
        
        question_type_full = ""
        if question_type == "choice":
            question_type_full = "Use multiple choice questions only."
        if question_type == "short":
            question_type_full = "Use short answer questions only."
        if question_type == "both":
            question_type_full = "Use both multiple choice and short answer questions."
            
        target_grade_full = ""
        if targetGrade == "none":
            target_grade_full = ""
        else:
            target_grade_full = f"This quiz should be made for a {targetGrade} grade class."
        
        if ctx:
            model = ChatOpenAI(model_name="gpt-3.5-turbo-0125", temperature=0.8)
            filled_template = QUIZ_GENERATE_TEMPLATE.replace("{system}", QUIZ_SYSTEM_PROMPT)\
                                                .replace("{questionType}", question_type_full)\
                                                .replace("{prompt}", user_prompt)\
                                                .replace("{targetGrade}", target_grade_full)
            prompt = ChatPromptTemplate.from_template(template=filled_template)
            retriever = vectorstore.as_retriever(search_type="similarity_score_threshold", search_kwargs={"score_threshold": 0.5, "k": 4})
        
            chain = (
                {"document": retriever, "ctx": RunnablePassthrough()}
                | prompt
                | model
                | StrOutputParser()
            )
            response = chain.invoke(ctx)
            
            file_name = generate_quiz_from_XML(response, username)
            return JsonResponse({'filename': file_name, 'response': response, 'file_text': text, 'context': ctx})
        else:
            return JsonResponse({'response': 'No context specified'})
    return HttpResponse("Listening for requests...")
    
def serve_file(request, file_name):
    file_path = os.path.join(GENERATEDCONTENT_DIRECTORY, file_name)
    
    if os.path.exists(file_path) and os.path.isfile(file_path):
        download = request.GET.get('download', 'false').lower() in ['true', '1', 't']
        return FileResponse(open(file_path, 'rb'), as_attachment=download)
    else:
        raise Http404("The requested file does not exist")
    
@csrf_exempt
def regenerate_presentation(request):
    if (request.method == 'POST'):
        load_dotenv()
        document_text = request.POST.get('documentText')
        user_prompt = request.POST.get('prompt')
        filename = request.POST.get('filename')
        original_string = request.POST.get("originalString")
        bgcolor = request.POST.get('backgroundColor')
        fonttype = request.POST.get('fontType')
        fontcolor = request.POST.get('fontColor')
        ctx = request.POST.get('ctx')
        username = request.POST.get('username')
        
        # Split text into chunks
        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
    
        chunks = text_splitter.split_text(document_text)
        embeddings = OpenAIEmbeddings()
		
        vectorstore = FAISS.from_texts(chunks, embeddings)
        
        if user_prompt:
            model = ChatOpenAI(model_name="gpt-3.5-turbo-0125", temperature=0.8)
            filled_template = PRESENTATION_REGENERATE_TEMPLATE.replace("{originalString}", original_string)\
                                                 .replace("{system}", PRESENTATION_SYSTEM_PROMPT)\
                                                 .replace("{ctx}", ctx)
            prompt = ChatPromptTemplate.from_template(template=filled_template)
            retriever = vectorstore.as_retriever(search_type="similarity_score_threshold", search_kwargs={"score_threshold": 0.5, "k": 4})

            chain = (
                {"document": retriever, "prompt": RunnablePassthrough()}
                | prompt
                | model
                | StrOutputParser()
            )
            
            response = chain.invoke(user_prompt)
            
            apply_bgcolor = WHITE
            if bgcolor == 'black':
                apply_bgcolor = BLACK
            elif bgcolor == 'white':
                apply_bgcolor = WHITE
            elif bgcolor == 'lightblue':
                apply_bgcolor = LIGHTBLUE
            elif bgcolor == 'cream':
                apply_bgcolor = CREAM
            elif bgcolor == 'grey':
                apply_bgcolor = GREY

            apply_fontcolor = BLACK
            if fontcolor == 'black':
                apply_fontcolor = BLACK
            elif fontcolor == 'white':
                apply_fontcolor = WHITE
            file_name = generate_slides_from_XML(response, apply_bgcolor, fonttype, apply_fontcolor, username)
            
            # Delete old iteration from filesystem
            file_path_pptx = os.path.join(GENERATEDCONTENT_DIRECTORY, filename)
            base, ext = os.path.splitext(file_path_pptx)
            file_path_pdf = base + ".pdf"
            if os.path.exists(file_path_pptx):
                os.remove(file_path_pptx)
                print(f"File at {file_path_pptx} removed.")
            else:
                print("File does not exist")
            if os.path.exists(file_path_pdf):
                os.remove(file_path_pdf)
                print(f"File at {file_path_pdf} removed.")
            else:
                print("File does not exist")
            
            return JsonResponse({'filename' : file_name, 'response': response, 'file_text': document_text, 'style': {'bg': bgcolor, 'fontcolor': fontcolor, 'fonttype': fonttype}})  
        else:
            return JsonResponse({'response': 'No prompt specified'})
    return HttpResponse("Listening for requests on presentation regenerate...")

@csrf_exempt
def regenerate_quiz(request):
    if (request.method == 'POST'):
        load_dotenv()
        document_text = request.POST.get('documentText')
        user_prompt = request.POST.get('prompt')
        filename = request.POST.get('filename')
        original_string = request.POST.get("originalString")
        question_type = request.POST.get("questionType")
        ctx = request.POST.get('ctx')
        username = request.POST.get('username')
        
        # Split text into chunks
        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
    
        chunks = text_splitter.split_text(document_text)
        embeddings = OpenAIEmbeddings()
		
        vectorstore = FAISS.from_texts(chunks, embeddings)
        
        question_type_full = ""
        if question_type == "choice":
            question_type_full = "Use multiple choice questions only."
        if question_type == "short":
            question_type_full = "Use short answer questions only."
        if question_type == "both":
            question_type_full = "Use both multiple choice and short answer questions."
        
        if user_prompt:
            model = ChatOpenAI(model_name="gpt-3.5-turbo-0125", temperature=0.8)
            filled_template = QUIZ_REGENERATE_TEMPLATE.replace("{system}", QUIZ_SYSTEM_PROMPT)\
                                                    .replace("{ctx}", ctx)\
                                                    .replace("{questionType}", question_type_full)\
                                                    .replace("{originalString}", original_string)
            prompt = ChatPromptTemplate.from_template(template=filled_template)
            retriever = vectorstore.as_retriever(search_type="similarity_score_threshold", search_kwargs={"score_threshold": 0.5, "k": 4})
            
            chain = (
                {"document": retriever, "prompt": RunnablePassthrough()}
                | prompt
                | model
                | StrOutputParser()
            )
            
            response = chain.invoke(user_prompt)
            file_name = generate_quiz_from_XML(response, username)
            
            # Delete old iteration from filesystem
            file_path_docx = os.path.join(GENERATEDCONTENT_DIRECTORY, filename)
            base, ext = os.path.splitext(file_path_docx)
            file_path_pdf = base + ".pdf"
            if os.path.exists(file_path_docx):
                os.remove(file_path_docx)
                print(f"File at {file_path_docx} removed.")
            else:
                print("File does not exist")
            if os.path.exists(file_path_pdf):
                os.remove(file_path_pdf)
                print(f"File at {file_path_pdf} removed.")
            else:
                print("File does not exist")
                
            return JsonResponse({'filename': file_name, 'response': response, 'file_text': document_text})
        else:
            return JsonResponse({'response': 'No prompt specified'})
    return HttpResponse("Listening for requests on quiz regenerate...")

def convert_to_pdf(file_path, output_dir):
    try:
        subprocess.run([
            'soffice', '--convert-to', 'pdf', '--outdir', output_dir, file_path
        ], check=True)
        print(f"Converted {file_path} to PDF successfully.")
    except subprocess.CalledProcessError as e:
        print(f"Failed to convert {file_path} to PDF. Error: {e}")